"""
create_presentation.py
======================
Generates the TwinPaths project PPTX presentation.

Usage:
    python create_presentation.py

Output:
    results/TwinPaths_presentation.pptx

Requires:
    pip install python-pptx
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── colour palette ────────────────────────────────────────────────────────────
NAVY       = RGBColor(0x1F, 0x38, 0x64)
BLUE       = RGBColor(0x2E, 0x74, 0xB5)
LIGHT_BLUE = RGBColor(0xBD, 0xD7, 0xEE)
ACCENT     = RGBColor(0x00, 0xB0, 0xF0)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
DARK       = RGBColor(0x1A, 0x1A, 0x1A)
MID_GRAY   = RGBColor(0x76, 0x76, 0x76)
LIGHT_GRAY = RGBColor(0xF0, 0xF2, 0xF5)
GREEN      = RGBColor(0x21, 0x7B, 0x45)
LIGHT_GREEN= RGBColor(0xE2, 0xEF, 0xDA)
RED        = RGBColor(0xC0, 0x00, 0x00)
ORANGE     = RGBColor(0xD4, 0x70, 0x1C)

# ── presentation setup ────────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]   # completely blank layout


# ─── low-level primitives ─────────────────────────────────────────────────────

def new_slide():
    return prs.slides.add_slide(BLANK)


def rect(sl, x, y, w, h, fill=None, line_color=None):
    """Add a rectangle (x,y,w,h in inches)."""
    s = sl.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    if fill:
        s.fill.solid()
        s.fill.fore_color.rgb = fill
    else:
        s.fill.background()
    if line_color:
        s.line.color.rgb = line_color
    else:
        s.line.fill.background()
    return s


def txt(sl, text, x, y, w, h,
        size=17, bold=False, italic=False,
        color=DARK, align=PP_ALIGN.LEFT):
    """Add a text box."""
    tb = sl.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tb.word_wrap = True
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    return tb


def bullets(sl, items, x, y, w, h=5.5, size=17, color=DARK):
    """
    items: list of str (level-0) or (str, level).
    Levels 0,1 are supported.
    """
    tb = sl.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tb.word_wrap = True
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        level = 0
        if isinstance(item, tuple):
            text, level = item
        else:
            text = item
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_before = Pt(5 if level == 0 else 2)
        indent = "    " * level
        bullet = "▸  " if level == 0 else "–  "
        r = p.add_run()
        r.text = indent + bullet + text
        r.font.size = Pt(size - level * 2)
        r.font.color.rgb = color
    return tb


def formula_box(sl, lines, x, y, w, size=14):
    """Light-gray box for formulas / code; returns bottom y."""
    line_h = 0.36
    h = line_h * len(lines) + 0.18
    rect(sl, x, y, w, h, fill=LIGHT_GRAY, line_color=MID_GRAY)
    for i, line in enumerate(lines):
        txt(sl, line, x + 0.15, y + 0.07 + i * line_h, w - 0.3, line_h + 0.05,
            size=size, color=NAVY)
    return y + h   # bottom edge


def header(sl, title, subtitle=None):
    """Standard slide header: navy band with white title."""
    band_h = 1.15 if subtitle else 1.05
    rect(sl, 0, 0, 13.33, band_h, fill=NAVY)
    txt(sl, title, 0.3, 0.06, 12.7, 0.8, size=28, bold=True, color=WHITE)
    if subtitle:
        txt(sl, subtitle, 0.3, 0.8, 12.7, 0.38, size=13, color=LIGHT_BLUE)


def tag(sl, label, x, y, fill=BLUE):
    """Small coloured pill label."""
    w = max(len(label) * 0.115 + 0.35, 1.1)
    rect(sl, x, y, w, 0.38, fill=fill)
    txt(sl, label, x + 0.08, y + 0.04, w - 0.16, 0.32,
        size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)


def bottom_note(sl, text, y=6.55, fill=LIGHT_GRAY, color=NAVY, italic=False):
    rect(sl, 0.3, y, 12.73, 0.7, fill=fill, line_color=MID_GRAY)
    txt(sl, text, 0.45, y + 0.08, 12.43, 0.55, size=13, color=color, italic=italic)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ═══════════════════════════════════════════════════════════════════════════════
sl = new_slide()
rect(sl, 0, 0, 13.33, 7.5, fill=NAVY)
rect(sl, 0, 4.75, 13.33, 0.08, fill=BLUE)
rect(sl, 0, 4.83, 13.33, 0.04, fill=LIGHT_BLUE)

txt(sl, "Dual-Path Trees",
    0.5, 1.0, 12.33, 1.4, size=54, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
txt(sl, "Three Algorithms for Survivable Network Design",
    0.5, 2.5, 12.33, 0.8, size=26, color=LIGHT_BLUE, align=PP_ALIGN.CENTER)
txt(sl, "DPGC Heuristic   ·   Matroid Intersection   ·   Brute-Force Ground Truth",
    0.5, 3.4, 12.33, 0.55, size=16, color=ACCENT, align=PP_ALIGN.CENTER)
txt(sl, "[Author]   |   [Course]   |   Technion, 2026",
    0.5, 5.2, 12.33, 0.5, size=15, color=MID_GRAY, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — The Problem
# ═══════════════════════════════════════════════════════════════════════════════
sl = new_slide()
header(sl, "The Problem: Survivable Network Design",
       "How do we design networks that survive edge failures?")

bullets(sl, [
    "Networks (internet, telecom, power grids) must stay connected under failures",
    "A single edge failure can disconnect nodes — unacceptable for critical infrastructure",
    "k-survivability: guarantee k edge-disjoint paths between terminals s and t",
    "Here: k = 2  →  two independent routes always exist",
    "We also require the subgraph to span all nodes (not just connect s to t)",
], x=0.4, y=1.3, w=8.3, size=18)

rect(sl, 9.0, 1.3, 4.0, 3.0, fill=LIGHT_BLUE, line_color=BLUE)
txt(sl, "The Goal", 9.15, 1.4, 3.7, 0.4, size=14, bold=True, color=NAVY)
txt(sl,
    "Find the minimum-cost spanning subgraph H ⊆ G that contains two node-disjoint s–t paths.",
    9.15, 1.85, 3.7, 2.2, size=15, color=DARK)

bottom_note(sl,
    "★  Known as the Dual-Path Tree (DPT) or Dual-Path Spanning Tree (DPST) problem. "
    "Reference: Balakrishnan et al. (1998, 2004).")


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Formal Definition
# ═══════════════════════════════════════════════════════════════════════════════
sl = new_slide()
header(sl, "Formal Problem Definition")

txt(sl, "Input", 0.4, 1.2, 1.5, 0.4, size=15, bold=True, color=BLUE)
bot = formula_box(sl, [
    "G = (V, E, w)   — undirected weighted graph",
    "s, t ∈ V        — source and destination terminals",
], x=0.4, y=1.65, w=12.5, size=15)

txt(sl, "Output", 0.4, bot + 0.2, 1.5, 0.4, size=15, bold=True, color=BLUE)
bot = formula_box(sl, [
    "H ⊆ G   spanning all nodes  (H includes every v ∈ V)",
    "H contains two internally node-disjoint s–t paths",
    "w(H) = Σ_{e ∈ H} w(e)  is minimised",
], x=0.4, y=bot + 0.6, w=12.5, size=15)

txt(sl, "Key observations", 0.4, bot + 0.25, 4.0, 0.4, size=15, bold=True, color=NAVY)
bullets(sl, [
    "H must connect ALL nodes — not just s and t",
    "Node-disjoint paths ⟹ edge-disjoint  (node-disjoint is the stronger condition)",
    "NP-hard in general; polynomial when the triangle inequality holds on edge weights",
    "Triangle inequality: w(u,v) ≤ w(u,x) + w(x,v)  for all u,v,x — satisfied by Euclidean distances",
], x=0.4, y=bot + 0.7, w=12.5, size=15)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Illustrated Example
# ═══════════════════════════════════════════════════════════════════════════════
sl = new_slide()
header(sl, "Example: What Does a DPT Solution Look Like?")

bullets(sl, [
    "Graph G with 6 nodes, s = 1, t = 6",
    "Two node-disjoint paths found:",
    ("Path 1  (red):  1 → 2 → 5 → 6", 1),
    ("Path 2  (blue): 1 → 3 → 4 → 6", 1),
    "All 6 nodes are covered — it's a spanning subgraph",
    "Total cost = sum of all selected edge weights",
    "Any unvisited node would be attached via cheapest available edge",
], x=0.4, y=1.25, w=5.9, size=17)

rect(sl, 6.7, 1.2, 6.3, 5.8, fill=LIGHT_GRAY, line_color=MID_GRAY)
txt(sl, "DPT solution structure", 6.9, 1.3, 5.9, 0.4, size=14, bold=True, color=NAVY)
txt(sl, """\
                 s  (node 1)
                /            \\
         [P1] /               \\ [P2]
             /                 \\
         node 2             node 3
             |                 |
         node 5             node 4
               \\             /
           [P1] \\           / [P2]
                 \\         /
                  t  (node 6)


  ── Path 1  (red)      ── Path 2  (blue)
  ·· Additional spanning edges (none here)\
""",
    6.85, 1.75, 6.05, 5.1, size=13, color=DARK)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — Algorithm 1: Brute-Force
# ═══════════════════════════════════════════════════════════════════════════════
sl = new_slide()
header(sl, "Algorithm 1: Brute-Force  (Ground Truth)")
tag(sl, "EXACT", 12.35, 0.33, fill=GREEN)

bullets(sl, [
    "Step 1 — Path enumeration:  DFS to list ALL simple s–t paths in G",
    "Step 2 — Pair search:  try every ordered pair (P₁, P₂) of paths",
    ("Reject if they share an intermediate node or an edge", 1),
    "Step 3 — Greedy attachment:  for remaining nodes, repeatedly add the",
    ("cheapest edge connecting an unvisited node to the current subgraph", 1),
    "Step 4 — Track the best (P₁, P₂) + attachment set by total cost",
    "Step 5 — Return optimal solution and its cost",
], x=0.4, y=1.2, w=8.6, size=17)

rect(sl, 9.35, 1.2, 3.65, 3.6, fill=RGBColor(0xFF, 0xF0, 0xF0), line_color=RED)
txt(sl, "Complexity", 9.5, 1.3, 3.3, 0.4, size=14, bold=True, color=RED)
txt(sl,
    "Path count:   up to O(n!)\n\n"
    "Pair loop:    O(paths²)\n\n"
    "Attachment:   O(n²) per pair\n\n"
    "⚠  Exponential in n\n\n"
    "Practical only for  n ≤ 9",
    9.5, 1.75, 3.3, 2.8, size=14, color=DARK)

bottom_note(sl,
    "✔  Purpose: produces the proven-optimal answer on small graphs "
    "to validate the two polynomial algorithms.")


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — DPGC Heuristic Overview
# ═══════════════════════════════════════════════════════════════════════════════
sl = new_slide()
header(sl, "Algorithm 2: DPGC Heuristic",
       "Balakrishnan et al. (2004) — polynomial time, not always optimal")
tag(sl, "HEURISTIC", 11.15, 0.33, fill=ORANGE)

bullets(sl, [
    "Decomposes the problem into two tractable sub-problems",
    "Step 1 — Two min-cost edge-disjoint s–t paths  (min-cost flow, k=2)",
    "Step 2 — Contract all dual-path nodes  N₁  into a single supernode  C",
    "Step 3 — Metric closure of the contracted graph, then MST",
    "Step 4 — Map MST edges back to original graph edges",
    "Final solution:  E₁  (dual-path edges)  ∪  recovered MST edges",
], x=0.4, y=1.25, w=12.5, size=18)

# flow diagram
step_labels = [
    ("Min-Cost\nFlow (k=2)", BLUE),
    ("Contract\nG[N₁] → C", GREEN),
    ("MST on\nClosure", ORANGE),
    ("Recover\nOriginal\nEdges", NAVY),
]
for i, (label, c) in enumerate(step_labels):
    xp = 0.45 + i * 3.15
    rect(sl, xp, 5.1, 2.85, 1.2, fill=c)
    txt(sl, label, xp + 0.07, 5.15, 2.71, 1.1, size=15, bold=True,
        color=WHITE, align=PP_ALIGN.CENTER)
    if i < 3:
        txt(sl, "→", xp + 2.85, 5.45, 0.3, 0.5, size=22, bold=True,
            color=NAVY, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — DPGC Step-by-Step Detail
# ═══════════════════════════════════════════════════════════════════════════════
sl = new_slide()
header(sl, "DPGC — Step-by-Step Detail")

for y_start, step_color, step_title, step_items in [
    (1.2, BLUE,   "Step 1: Two Min-Cost Edge-Disjoint Paths", [
        "Build directed flow network: each edge → two arcs, capacity 1, cost = weight",
        "Set demand: supply k=2 at s, demand k=2 at t; solve with network_simplex",
        "Decompose flow into two s–t paths → edge set E₁, node set N₁",
    ]),
    (2.75, GREEN,  "Step 2: Contraction", [
        "Replace G[N₁] (subgraph induced by path nodes) with single node 'C'",
        "Edge (C, v): weight = min weight among edges from any node in N₁ to v",
    ]),
    (3.9, ORANGE, "Step 3: Metric Closure + MST", [
        "Replace edge weights in contracted graph H with shortest-path distances",
        "Run Kruskal's MST on this closure → connects all remaining nodes optimally",
    ]),
    (4.95, NAVY,  "Step 4: Edge Recovery", [
        "Trace each MST edge back via stored shortest path → original edges in G",
        "Edges touching 'C' mapped back to the original node in N₁",
        "Final:  E₁ ∪ recovered_edges  (all on original node labels)",
    ]),
]:
    txt(sl, step_title, 0.4, y_start, 12.5, 0.4, size=15, bold=True, color=step_color)
    for k, item in enumerate(step_items):
        txt(sl, "    ▸  " + item, 0.4, y_start + 0.4 + k * 0.33, 12.5, 0.35, size=13, color=DARK)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — What is a Matroid?
# ═══════════════════════════════════════════════════════════════════════════════
sl = new_slide()
header(sl, "What is a Matroid?",
       "A combinatorial abstraction of linear independence")

txt(sl, "Formal Definition", 0.4, 1.2, 5.0, 0.4, size=16, bold=True, color=NAVY)
bot = formula_box(sl, [
    "A matroid  M = (E, I)  where:",
    "  E  —  finite ground set  (e.g. edges of a graph)",
    "  I  —  family of subsets of E  called independent sets",
], x=0.4, y=1.65, w=6.4, size=15)

txt(sl, "Three Axioms", 0.4, bot + 0.15, 3.0, 0.4, size=16, bold=True, color=NAVY)
formula_box(sl, [
    "(I1)  ∅ ∈ I                                          (empty set is independent)",
    "(I2)  A ∈ I, B ⊆ A   ⟹   B ∈ I                   (hereditary / downward-closed)",
    "(I3)  A, B ∈ I, |A| < |B|   ⟹   ∃ e ∈ B\\A : A∪{e} ∈ I    (augmentation)",
], x=0.4, y=bot + 0.6, w=12.5, size=14)

rect(sl, 6.9, 1.2, 6.1, 2.7, fill=LIGHT_BLUE, line_color=BLUE)
txt(sl, "Key vocabulary", 7.05, 1.3, 5.8, 0.4, size=14, bold=True, color=NAVY)
bullets(sl, [
    "Base — maximal independent set  (all bases have equal size)",
    "Circuit — minimal dependent set  (minimal set not in I)",
    "Rank  r(S) — size of largest independent subset of S",
    "Greedy algorithm finds a maximum-weight base of any matroid",
], x=7.0, y=1.78, w=5.85, size=13.5, color=DARK)

bottom_note(sl,
    "Intuition: Matroids capture exactly those independence systems "
    "where the greedy algorithm gives a globally optimal solution.",
    italic=True)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — Matroid Examples
# ═══════════════════════════════════════════════════════════════════════════════
sl = new_slide()
header(sl, "Matroid Examples")

for col_x, (title, blist) in enumerate([
    ("Graphic Matroid", [
        "E = edges of graph G",
        "I = acyclic subsets (forests)",
        "Bases = spanning trees",
        "Circuit = simple cycle",
        "Greedy = Kruskal's MST",
    ]),
    ("Uniform Matroid  Uₖ,ₙ", [
        "E = {1, 2, …, n}",
        "I = all subsets of size ≤ k",
        "Bases = all k-element subsets",
        "Circuit = any (k+1)-element set",
        "Greedy = select top-k elements",
    ]),
    ("Partition Matroid", [
        "E partitioned into blocks B₁…Bₘ",
        "I = sets with ≤ dᵢ elements from Bᵢ",
        "Models per-group quotas",
        "Circuit = exceeds quota in one block",
        "Greedy = local optimum per block",
    ]),
]):
    x0 = 0.3 + col_x * 4.35
    rect(sl, x0, 1.2, 4.1, 4.0, fill=LIGHT_GRAY, line_color=MID_GRAY)
    txt(sl, title, x0 + 0.15, 1.3, 3.8, 0.45, size=14, bold=True, color=NAVY)
    for k, item in enumerate(blist):
        txt(sl, "▸  " + item, x0 + 0.2, 1.82 + k * 0.58, 3.7, 0.55, size=13.5, color=DARK)

bottom_note(sl,
    "The DPT algorithm uses a non-standard matroid called the "
    "Q-Restricted 1-Tree Matroid — explained on the next slide.",
    fill=LIGHT_BLUE, color=NAVY)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — Q-Restricted 1-Tree Matroid
# ═══════════════════════════════════════════════════════════════════════════════
sl = new_slide()
header(sl, "The Q-Restricted 1-Tree Matroid",
       "The key building block of the DPT matroid algorithm")

txt(sl, "Building up the definition:", 0.4, 1.2, 8.0, 0.4, size=16, bold=True, color=NAVY)
bullets(sl, [
    "Spanning tree  —  acyclic subgraph connecting all n nodes  (exactly n-1 edges)",
    "1-tree          —  spanning tree + one extra edge  (exactly n edges, one cycle)",
    "Q-restricted 1-tree  —  a 1-tree where the unique cycle must pass through node q",
], x=0.4, y=1.65, w=12.5, size=17)

txt(sl, "As a matroid  M_q = (E, I_q):", 0.4, 3.05, 8.0, 0.4, size=16, bold=True, color=NAVY)
formula_box(sl, [
    "Ground set:     E = edges of G",
    "Independent:    I_q = { F ⊆ E  |  F is a forest,",
    "                         OR  F is a forest + exactly one cycle and q lies on that cycle }",
    "Circuit of e:   minimal set in I_q ∪ {e} that is dependent in M_q",
], x=0.4, y=3.5, w=12.5, size=14)

rect(sl, 0.4, 5.5, 12.5, 1.2, fill=LIGHT_BLUE, line_color=BLUE)
txt(sl, "Why this matroid?", 0.55, 5.58, 3.5, 0.4, size=14, bold=True, color=NAVY)
txt(sl,
    "A q-restricted 1-tree is a spanning structure that encodes the requirement "
    "that node q participates in any cycle. "
    "Intersecting M_s with M_t forces BOTH s and t to be on a cycle — "
    "which is exactly the condition for two node-disjoint s–t paths.",
    0.55, 5.98, 12.2, 0.65, size=13.5, italic=True, color=DARK)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — Matroid Intersection (Concept)
# ═══════════════════════════════════════════════════════════════════════════════
sl = new_slide()
header(sl, "Matroid Intersection",
       "Finding the maximum common independent set of two matroids")

txt(sl, "The Problem", 0.4, 1.2, 4.0, 0.4, size=16, bold=True, color=NAVY)
bot = formula_box(sl, [
    "Given  M₁ = (E, I₁)  and  M₂ = (E, I₂)  on the same ground set E,",
    "find   I* ∈ I₁ ∩ I₂   of maximum total weight   Σ_{e∈I*} w(e).",
], x=0.4, y=1.65, w=12.5, size=15)

bullets(sl, [
    "2-matroid intersection is solvable in polynomial time",
    ("Unlike 3-matroid intersection, which is NP-hard", 1),
    "Many classical combinatorial problems are special cases:",
    ("Bipartite matching  (partition matroid ∩ partition matroid)", 1),
    ("Arborescences, common bases, colourful spanning trees, …", 1),
    "Weighted version: maximise Σ w(e) — solved by augmenting-path algorithm (Lawler 1975)",
    "We use the weighted version, but minimise cost: set  w(e) = −weight(e)",
], x=0.4, y=bot + 0.25, w=12.5, size=17)

bottom_note(sl,
    "Key result (Lawler 1975): Weighted matroid intersection runs in "
    "O(r² · T_oracle) where r = |I*| and T_oracle is the time for a circuit query.",
    italic=True)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — The Exchange Graph & Algorithm
# ═══════════════════════════════════════════════════════════════════════════════
sl = new_slide()
header(sl, "Weighted Matroid Intersection — The Algorithm",
       "Augmenting paths in an exchange graph (Lawler / cplib port)")

txt(sl, "Exchange graph  D_I  (constructed from the current solution I):",
    0.4, 1.2, 12.5, 0.4, size=15, bold=True, color=NAVY)
formula_box(sl, [
    "Nodes:   each element e ∈ E  plus virtual source s* and sink t*",
    "Arcs:",
    "  e ∉ I, adding e is free in M₁ (no M₁-circuit)   →  arc  e → t*   cost 0",
    "  e ∉ I, f ∈ circuit_M₁(e) \\ {e}                  →  arc  e → f    cost −w(f)+1",
    "  e ∉ I, adding e is free in M₂                   →  arc  s* → e   cost  w(e)+1",
    "  e ∉ I, f ∈ circuit_M₂(e) \\ {e}                  →  arc  f → e    cost  w(e)+1",
], x=0.4, y=1.65, w=12.5, size=13)

txt(sl, "Augmentation loop:", 0.4, 4.15, 5.0, 0.4, size=15, bold=True, color=NAVY)
bullets(sl, [
    "Find the shortest path  s* → t*  in D_I  using Bellman-Ford  (handles negative weights)",
    "Augment:  toggle I[e] for every element e on the path  (symmetric difference)",
    "Rebuild D_I from the new I and repeat",
    "Terminate when no augmenting path exists  →  I is maximum-weight",
], x=0.4, y=4.6, w=12.5, size=16)

bottom_note(sl,
    "Each augmentation strictly increases total weight → terminates in O(|E|) steps. "
    "Negative-weight edges in D_I are the reason Bellman-Ford is used instead of Dijkstra.",
    fill=LIGHT_BLUE, color=NAVY)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 13 — Algorithm 3: DPT via Matroid Intersection
# ═══════════════════════════════════════════════════════════════════════════════
sl = new_slide()
header(sl, "Algorithm 3: DPT via Matroid Intersection",
       "Exact algorithm under the triangle inequality")
tag(sl, "EXACT*", 12.0, 0.33, fill=GREEN)

txt(sl, "The key insight:", 0.4, 1.2, 8.0, 0.4, size=16, bold=True, color=NAVY)
bot = formula_box(sl, [
    "A spanning subgraph H has two node-disjoint s–t paths",
    "     ⟺     H ∈ I(M_s) ∩ I(M_t)",
    "where  M_s = Q-restricted 1-tree matroid anchored at  q = s",
    "       M_t = Q-restricted 1-tree matroid anchored at  q = t",
], x=0.4, y=1.65, w=12.5, size=15)

txt(sl, "Algorithm:", 0.4, bot + 0.2, 4.0, 0.4, size=16, bold=True, color=NAVY)
bullets(sl, [
    "Build M_s and M_t on the edge set of G  (shared ground set E)",
    "Run weighted matroid intersection  maximising  −Σ w(e)  (i.e. minimising cost)",
    "I* ∈ I(M_s) ∩ I(M_t) gives the optimal DPT",
    "Extract two disjoint paths from the solution edge set",
], x=0.4, y=bot + 0.65, w=12.5, size=17)

rect(sl, 0.4, 5.6, 12.5, 1.1, fill=LIGHT_GREEN, line_color=GREEN)
txt(sl, "Why it is exact:", 0.55, 5.7, 3.5, 0.4, size=14, bold=True, color=GREEN)
txt(sl,
    "Under the triangle inequality the maximum-weight element of I(M_s) ∩ I(M_t) "
    "corresponds to a minimum-cost spanning subgraph with two node-disjoint s–t paths. "
    "Triangle inequality ensures no shortcut artificially reduces cost, "
    "so the matroid formulation is tight.",
    0.55, 6.07, 12.2, 0.58, size=13, italic=True, color=DARK)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 14 — The Triangle Inequality — Why It Matters
# ═══════════════════════════════════════════════════════════════════════════════
sl = new_slide()
header(sl, "Why the Triangle Inequality Matters")

bullets(sl, [
    "Triangle inequality: for every edge (u,v) and every node x,",
    ("w(u,v)  ≤  w(u,x) + w(x,v)", 1),
    "Satisfied automatically by Euclidean (geographic) distances",
    "If violated, the matroid algorithm may return a suboptimal or infeasible result",
    "Remedy: replace G with its metric closure  (edge weights = shortest-path distances)",
    ("Metric closure always satisfies the triangle inequality by construction", 1),
    ("May increase edge count to O(n²) but guarantees algorithm correctness", 1),
    "In this project: generators produce Euclidean graphs → triangle inequality is free",
], x=0.4, y=1.25, w=8.5, size=17)

rect(sl, 9.1, 1.2, 3.9, 3.8, fill=LIGHT_GRAY, line_color=MID_GRAY)
txt(sl, "Our graph generators", 9.25, 1.3, 3.6, 0.4, size=13, bold=True, color=NAVY)
bullets(sl, [
    "Euclidean random graph",
    "Complete Euclidean graph",
    "Grid graph (unit weights)",
], x=9.25, y=1.75, w=3.6, size=13, color=DARK)
txt(sl, "→ all satisfy\ntriangle inequality",
    9.25, 3.0, 3.6, 0.8, size=14, bold=True, color=GREEN, align=PP_ALIGN.CENTER)

bottom_note(sl,
    "The DPGC heuristic does not require the triangle inequality — "
    "it works on any 2-edge-connected graph.")


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 15 — Comparison: Solution Quality
# ═══════════════════════════════════════════════════════════════════════════════
sl = new_slide()
header(sl, "Comparison: Solution Quality")

# Table header
rect(sl, 0.3, 1.2, 12.73, 0.55, fill=NAVY)
for label, xp, wc in [
    ("Graph",               0.35, 2.9),
    ("Brute-Force (opt.)",  3.3,  3.1),
    ("Matroid DPT",         6.5,  3.0),
    ("DPGC Heuristic",      9.6,  3.2),
]:
    txt(sl, label, xp, 1.25, wc, 0.4, size=13, bold=True, color=WHITE)

rows = [
    ("random n=5  seed=0", "8.24",  "8.24",  "8.24  ✔",  True),
    ("random n=6  seed=1", "12.57", "12.57", "13.41  ✗", False),
    ("random n=7  seed=2", "15.03", "15.03", "15.03  ✔", True),
    ("random n=8  seed=0", "19.88", "19.88", "21.14  ✗", False),
    ("complete n=5",        "31.22", "31.22", "31.22  ✔", True),
    ("grid 3×3",            "6.00",  "6.00",  "6.00  ✔",  True),
]
alt = [WHITE, LIGHT_GRAY]
for i, (graph, bf, mat, dpgc, ok) in enumerate(rows):
    yp = 1.78 + i * 0.57
    rect(sl, 0.3, yp, 12.73, 0.57, fill=alt[i % 2], line_color=MID_GRAY)
    txt(sl, graph, 0.35, yp + 0.1, 2.9, 0.38, size=13, color=DARK)
    txt(sl, bf,    3.3,  yp + 0.1, 3.1, 0.38, size=13, color=DARK)
    txt(sl, mat,   6.5,  yp + 0.1, 3.0, 0.38, size=13, color=DARK)
    txt(sl, dpgc,  9.6,  yp + 0.1, 3.2, 0.38, size=13, color=(GREEN if ok else RED))

bottom_note(sl,
    "Note: costs above are illustrative. "
    "Run  compare_algorithms.py  on your graphs for real numbers.",
    italic=True, color=MID_GRAY)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 16 — Comparison: Runtime
# ═══════════════════════════════════════════════════════════════════════════════
sl = new_slide()
header(sl, "Comparison: Runtime vs. Graph Size",
       "Mean over random Euclidean graphs per size  (run: python -m experiments.benchmark_statistical)")

plot_path = os.path.join("results", "benchmark_runtimes.png")
if os.path.exists(plot_path):
    sl.shapes.add_picture(plot_path, Inches(1.8), Inches(1.3), Inches(9.7), Inches(4.7))
else:
    rect(sl, 1.8, 1.3, 9.7, 4.7, fill=LIGHT_GRAY, line_color=MID_GRAY)
    txt(sl, "[benchmark_runtimes.png]\n\nRun  python -m experiments.benchmark_statistical  to generate this plot.",
        2.3, 3.0, 8.7, 1.4, size=16, color=MID_GRAY, align=PP_ALIGN.CENTER)

bullets(sl, [
    "DPGC Heuristic:  O(n³ + n² log n)  — min-cost flow + MST",
    "Matroid DPT:     O(n² · |E|²)      — Bellman-Ford per augmentation",
    "Brute-Force:     O(n!)              — exponential, reference only",
], x=0.3, y=6.2, w=12.73, size=14)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 17 — Algorithm Summary Table
# ═══════════════════════════════════════════════════════════════════════════════
sl = new_slide()
header(sl, "Algorithm Summary")

rect(sl, 0.3, 1.2, 12.73, 0.55, fill=NAVY)
for label, xp, wc in [
    ("",                  0.35, 2.1),
    ("Brute-Force",       2.5,  3.2),
    ("DPGC Heuristic",    5.8,  3.4),
    ("Matroid DPT",       9.3,  3.5),
]:
    txt(sl, label, xp, 1.25, wc, 0.4, size=13, bold=True, color=WHITE)

table_rows = [
    ("Optimality",     "Optimal",          "Heuristic — may miss", "Optimal* (metric)"),
    ("Complexity",     "O(n!)  exp.",       "Polynomial — fast",     "Polynomial — moderate"),
    ("Max graph size", "n ≤ 9",            "Any size",              "~n ≤ 200"),
    ("Triangle ineq.", "Not required",      "Not required",          "Required"),
    ("Best use case",  "Correctness check", "Large prod. graphs",    "Exact on metric graphs"),
]
alt = [WHITE, LIGHT_GRAY]
for i, (prop, bf, dpgc, mat) in enumerate(table_rows):
    yp = 1.78 + i * 0.77
    rect(sl, 0.3, yp, 12.73, 0.77, fill=alt[i % 2], line_color=MID_GRAY)
    rect(sl, 0.3, yp, 2.15, 0.77, fill=LIGHT_BLUE, line_color=MID_GRAY)
    txt(sl, prop, 0.35, yp + 0.2, 2.05, 0.4, size=13, bold=True, color=NAVY)
    for label, xp, wc in [
        (bf,   2.5, 3.2),
        (dpgc, 5.8, 3.4),
        (mat,  9.3, 3.5),
    ]:
        txt(sl, label, xp, yp + 0.2, wc, 0.4, size=13, color=DARK)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 18 — Conclusions
# ═══════════════════════════════════════════════════════════════════════════════
sl = new_slide()
header(sl, "Conclusions")

bullets(sl, [
    "The DPT problem is a meaningful model for survivable network design with k=2",
    "Three algorithms cover the full speed–optimality trade-off:",
    ("Brute-force:       exact, exponential — use only as a correctness reference", 1),
    ("DPGC heuristic:   polynomial and fast — sometimes misses the optimum", 1),
    ("Matroid DPT:      polynomial and exact — requires triangle inequality", 1),
    "Matroids elegantly encode the structural constraint:",
    ("Two node-disjoint s–t paths  ⟺  I ∈ I(M_s) ∩ I(M_t)", 1),
    "Matroid intersection is a powerful technique applicable far beyond this problem",
    "Triangle inequality is not just a convenience — it ensures the matroid formulation is tight",
], x=0.4, y=1.3, w=12.5, size=17)

rect(sl, 0.3, 5.85, 12.73, 1.35, fill=NAVY)
txt(sl, "Thank You",
    0.3, 5.88, 12.73, 0.85, size=36, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
txt(sl, "Code:  github.com/[repo]   |   References: Balakrishnan et al. (1998, 2004) · Lawler (1975)",
    0.3, 6.7, 12.73, 0.4, size=12, color=LIGHT_BLUE, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
# Save
# ═══════════════════════════════════════════════════════════════════════════════
os.makedirs("results", exist_ok=True)
out_path = os.path.join(os.path.dirname(os.path.abspath(__file__)),
                        "results", "TwinPaths_presentation.pptx")
prs.save(out_path)
print(f"Saved → {out_path}")
